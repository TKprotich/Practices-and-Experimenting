from pptx import Presentation 
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.util import Pt

from datetime import datetime
import os

import pandas as pd


from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

folder = r"C:\Users\user\Downloads"
now = datetime.now()
today = now.strftime("%d %B")
today = "02 November Per Doctors"
columns = ["DOCTOR", 'No/P', 'Biochemistry', 'Hormone', 'hematology', 'serology', 'Urine', 'Stool', 'Coag', 'Semen', 'Swab', 'culture']
bigdata = pd.DataFrame(columns)
#search file
rawdata = r"C:\Users\user\Documents\2020-Kalkaal Speciality Hospital\Abdullah\Daily Data\31October Per Doctors.xlsx"
try:
    fh = pd.read_excel(rawdata, sheet_name = 'Sheet1 (2)')
except Exception as e:
    print(e, " in ", name)
#data cleaning
fh.drop_duplicates(inplace = True)
fh = fh.drop([1])
fh = fh.dropna()
fh = fh.set_index("DOCTOR")
print(fh)      
fh.drop_duplicates(inplace = True)
#fh = fh.T
prs = Presentation(r"C:\Users\user\Documents\2020-Kalkaal Speciality Hospital\Abdullah\New Microsoft PowerPoint Presentation.pptx")
_slide_layout = prs.slide_layouts[5]
_section_layout = prs.slide_layouts[2]
def onegraph(fh):
    fh1 = fh
        ###
    slide = prs.slides.add_slide(_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = f"Total Requests Per Doctor {today}"
    chart_data = CategoryChartData()
    chart_data.categories = [l for l in fh1.index]
    chart_data.add_series('Series 1', (k for k in fh[fh.columns[0]]))
    # add chart to slide --------------------
    x, y, cx, cy = Inches(1), Inches(3), Inches(11), Inches(4)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart
    #Categories
    category_axis = chart.category_axis
    category_axis.tick_labels.font.italic = True
    category_axis.tick_labels.font.size = Pt(15)

    chart.chart_title.has_text_frame=True
    chart.chart_title.text_frame.text= 'Total Patients'
    #Value axis
    value_axis = chart.value_axis
    #value_axis.maximum_scale = 50.0
    #Labels
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(17)
    data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
def presentation(fh):
    for i in range(len(fh.index)):
        slide = prs.slides.add_slide(_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = '{}'.format(fh.index[i])
        
        # define chart data ---------------------
        chart_data = CategoryChartData()
        chart_data.categories = [l for l in fh.columns[1:]]
        #chart_data.add_series('Series 1', (k for k in fh[i])) i for i in 
        chart_data.add_series('Series 1', (k for k in fh.iloc[i, 1:]))
        # add chart to slide --------------------
        x, y, cx, cy = Inches(1), Inches(3), Inches(11), Inches(4)
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )
        #print([l for l in fh.index[1:]])
        #print(k for k in fh[i][1:])
        chart = graphic_frame.chart
        #Categories
        category_axis = chart.category_axis
        category_axis.tick_labels.font.italic = True
        category_axis.tick_labels.font.size = Pt(15)

        chart.chart_title.has_text_frame=True
        chart.chart_title.text_frame.text= 'Total Patients {}'.format(fh['No/P'][i])
        #Value axis
        value_axis = chart.value_axis
        #value_axis.maximum_scale = 50.0
        
        #Labels
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(17)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
def presentation1(fh):
    onegraph(fh)
    for i in fh.columns:
        slide = prs.slides.add_slide(_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        
        title_shape.text = '{}'.format(i)

        # define chart data ---------------------
        chart_data = CategoryChartData()
        #chart_data.categories = [l for l in fh.index]
        chart_data.categories = [l for l in fh.index[1:]]
        
        #chart_data.add_series('Series 1', (k for k in fh[i]))
        chart_data.add_series('Series 1', (k for k in fh[i][1:]))
        # add chart to slide --------------------
        x, y, cx, cy = Inches(1), Inches(3), Inches(11), Inches(4)
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )
        #print([l for l in fh.index[1:]])
        #print(k for k in fh[i][1:])
        chart = graphic_frame.chart
        #Categories
        category_axis = chart.category_axis
        category_axis.tick_labels.font.italic = True
        category_axis.tick_labels.font.size = Pt(15)

        #chart.chart_title.has_text_frame=True
        #chart.chart_title.text_frame.text= 'Total Patients {}'.format(fh[i][0])
        #Value axis
        value_axis = chart.value_axis
        #value_axis.maximum_scale = 50.0
        
        #Labels
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(17)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    presentation(fh)
#presentation1()
def main():
    presentation1(fh)
    slide = prs.slides.add_slide(_section_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = f" End of Presentation on  Requests Per Doctor  for {today} By Abdullah Lab "
    prs.save("C:/Users/user/Documents/2020-Kalkaal Speciality Hospital/Abdullah/ {}.pptx".format(today))
if __name__ == "__main__":
    main()
