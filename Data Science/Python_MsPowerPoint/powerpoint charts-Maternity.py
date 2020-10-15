from pptx import Presentation 
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.util import Pt

from datetime import datetime
import os

import pandas as pd
import numpy as np


from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

folder = r"C:\Users\user\Downloads"
now = datetime.now()
today = now.strftime("%d %B")
today = "Inpatient Admission Report (14)"
for root, dirs, names in os.walk(folder):
    for name in names:
        if today.lower() in name.lower():
            rawdata = os.path.join(root, name)
prs = Presentation(r"C:\Users\user\Documents\2020-Kalkaal Speciality Hospital\Abdullah\Presentation1.pptx")
fh = pd.read_excel(rawdata, skiprows = 1)
#data cleaning
fh.set_index(fh.columns[0], inplace = True)
fh['Age'] = fh['A/S'].str.split(' ').str.get(0)
fh['Gender'] = fh['A/S'].str.split('/').str.get(1).str.capitalize()
fh.dropna(inplace = True)

#################### Counter
def counttext(column):
    column = pd.Series(column).dropna()
    column = list(chain(*column))
    counter = Counter(column)
    return counter

def dgdmydict( fh, analyzecolumn, filtercolumn, filteritem, DELIMETER):
    fh1 = fh[fh[filtercolumn] == filteritem]
    return counttext(fh1[analyzecolumn].str.lower().str.split(DELIMETER))

def analyzenow(fh, analyzecolumn, filtercolumn, filteritem):
    DELIMETER = ""
    newfh = fh[analyzecolumn].replace(' ', np.nan)
    newfh = fh[analyzecolumn].dropna()
    want = dgdmydict(fh, analyzecolumn, filtercolumn, filteritem, DELIMETER)
    thedict = dict(want.most_common())
    df = pd.DataFrame(want.most_common())
    df = df.rename(columns={0 : analyzecolumn, 1 : "Counts"})
    return df


#### Baby GENDER ###
gender = pd.DataFrame(fh.Gender.value_counts())

################## Generate Powerpoint
_slide_layout = prs.slide_layouts[5]
def presentation(fh):
    for i in fh.columns:
        slide = prs.slides.add_slide(_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        
        title_shape.text = '{}'.format(i)

        # define chart data ---------------------
        chart_data = CategoryChartData()
        chart_data.categories = [l for l in fh.index]
        
        chart_data.add_series('Series 1', (k for k in fh[i]))
        

        # add chart to slide --------------------
        x, y, cx, cy = Inches(1), Inches(3), Inches(11), Inches(4)
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )
        chart = graphic_frame.chart
        #Categories
        category_axis = chart.category_axis
        category_axis.tick_labels.font.italic = True
        category_axis.tick_labels.font.size = Pt(15)

        chart.chart_title.has_text_frame=True
        chart.chart_title.text_frame.text= '{}'.format(fh[i].name)
        #Value axis
        value_axis = chart.value_axis
        #value_axis.maximum_scale = 50.0
        
        #Labels
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(17)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

presentation(gender)
prs.save('chart-02.pptx')
